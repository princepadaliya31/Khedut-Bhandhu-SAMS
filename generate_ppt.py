from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

# ── colours ──────────────────────────────────────────────
BG     = RGBColor(0x0D, 0x2B, 0x1A)   # dark green
ACCENT = RGBColor(0x1A, 0x4A, 0x2E)   # mid green
GOLD   = RGBColor(0xD4, 0xAF, 0x37)   # gold
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY  = RGBColor(0xCC, 0xCC, 0xCC)

BLANK = prs.slide_layouts[6]           # completely blank layout

# ── helpers ──────────────────────────────────────────────
def add_bg(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def box(slide, l, t, w, h, bg=None, border=None):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.line.fill.background() if border is None else None
    if bg:
        shape.fill.solid(); shape.fill.fore_color.rgb = bg
    else:
        shape.fill.background()
    if border:
        shape.line.color.rgb = border
        shape.line.width = Pt(1.2)
    return shape

def txt(slide, text, l, t, w, h, size=18, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb

def header_bar(slide, title_text):
    box(slide, 0, 0, 13.33, 0.45, bg=RGBColor(0x0A, 0x20, 0x10))
    txt(slide, "Khedut Bandhu – SAMS", 0.2, 0.05, 6, 0.35, 11, False, GOLD)
    txt(slide, "Dharmsinh Desai University", 7, 0.05, 6, 0.35, 11, False, LGRAY, PP_ALIGN.RIGHT)
    txt(slide, title_text, 0.4, 0.55, 12, 0.6, 26, True, WHITE)
    box(slide, 0.4, 1.15, 2.5, 0.04, bg=GOLD)

def footer(slide, num):
    box(slide, 0, 7.1, 13.33, 0.4, bg=RGBColor(0x0A, 0x20, 0x10))
    txt(slide, f"Slide {num}", 12.5, 7.12, 0.7, 0.3, 10, False, LGRAY, PP_ALIGN.RIGHT)

def bullet_block(slide, items, l, t, w, h, gap=0.38):
    for i, (bold_part, rest) in enumerate(items):
        txb = slide.shapes.add_textbox(Inches(l), Inches(t + i*gap), Inches(w), Inches(gap+0.05))
        txb.word_wrap = True
        tf = txb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        if bold_part:
            r1 = p.add_run(); r1.text = bold_part + " "
            r1.font.bold = True; r1.font.size = Pt(13); r1.font.color.rgb = GOLD
        r2 = p.add_run(); r2.text = rest
        r2.font.size = Pt(13); r2.font.color.rgb = WHITE

# ════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_bg(sl)
box(sl, 0, 0, 13.33, 0.5, bg=RGBColor(0x0A,0x20,0x10))
# top header
txt(sl, "Department of Information Technology  |  Faculty of Technology",
    0, 0.08, 13.33, 0.4, 12, False, LGRAY, PP_ALIGN.CENTER)

# DDU crest placeholder circle
crest = sl.shapes.add_shape(9, Inches(5.7), Inches(0.7), Inches(1.9), Inches(1.9))
crest.fill.solid(); crest.fill.fore_color.rgb = GOLD
crest.line.color.rgb = WHITE; crest.line.width = Pt(2)
txt(sl, "DDU", 5.7, 1.1, 1.9, 0.7, 22, True, BG, PP_ALIGN.CENTER)

txt(sl, "Dharmsinh Desai University", 0, 2.7, 13.33, 0.55, 24, True, WHITE, PP_ALIGN.CENTER)
txt(sl, "Nadiad – 387001, Gujarat", 0, 3.22, 13.33, 0.4, 13, False, LGRAY, PP_ALIGN.CENTER)

box(sl, 1.5, 3.7, 10.33, 0.04, bg=GOLD)

txt(sl, "🌾  KHEDUT BANDHU", 0, 3.85, 13.33, 0.7, 32, True, GOLD, PP_ALIGN.CENTER)
txt(sl, "Smart Agriculture Management System  (SAMS)", 0, 4.5, 13.33, 0.5, 17, True, WHITE, PP_ALIGN.CENTER)
txt(sl, "B.Tech – Information Technology  |  Semester VI  |  Academic Year 2025–2026",
    0, 5.05, 13.33, 0.4, 12, False, LGRAY, PP_ALIGN.CENTER)

box(sl, 1.5, 5.55, 10.33, 0.04, bg=GOLD)

txt(sl, "Prepared By:", 1.5, 5.7, 4.5, 0.35, 13, True, GOLD)
txt(sl, "Padaliya Princekumar R.  (IT - 085)\nPatel Ansh K.  (IT - 091)", 1.5, 6.0, 5, 0.7, 13, False, WHITE)

txt(sl, "Guided By:", 7.5, 5.7, 4.5, 0.35, 13, True, GOLD)
txt(sl, "Prof. Sunil K. Vithlani\nDept. of Information Technology", 7.5, 6.0, 5, 0.7, 13, False, WHITE)

txt(sl, "April  2026", 0, 7.1, 13.33, 0.35, 11, False, LGRAY, PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════
# SLIDE 2 — Project Overview
# ════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_bg(sl)
header_bar(sl, "Project Overview")
footer(sl, 2)

items = [
    ("What is it?", "A full-stack Smart Agriculture Management System for Indian farmers, buyers & govt departments"),
    ("Platform Type:", "Web Application — MERN Stack + Python AI Service"),
    ("Target Users:", "Farmers  |  Buyers  |  Admin  |  Department Officials"),
    ("Project Type:", "B.Tech IT Semester VI — Academic Year 2025-26"),
    ("University:", "Dharmsinh Desai University, Nadiad — Dept. of Information Technology"),
    ("Guide:", "Prof. Sunil K. Vithlani, Dept. of IT, DDU"),
    ("Students:", "Padaliya Princekumar R. (IT-085)  &  Patel Ansh K. (IT-091)"),
    ("GitHub:", "github.com/apk3206/Khedut-Bandhu"),
]
bullet_block(sl, items, 0.5, 1.35, 12.3, 5.5, 0.6)

# ════════════════════════════════════════════════════════
# SLIDE 3 — Key Features
# ════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_bg(sl)
header_bar(sl, "Key Features")
footer(sl, 3)

features = [
    ("🔐", "Role-Based Auth", "Farmer / Buyer / Admin / Dept Admin\nOTP Email + Google OAuth + JWT"),
    ("🌿", "AI Crop Detection", "EfficientNetB3 CNN via FastAPI\nCLAHE preprocessing, 38 disease classes"),
    ("🛒", "E-Commerce Store", "Seeds, Pesticides, Tools\nCart, Checkout, UPI/Card/COD"),
    ("📊", "Market Analytics", "Live APMC crop prices\n6-month historical trend charts"),
    ("📋", "Complaint Mgmt", "Auto dept-routing by keywords\nEmail updates to farmer"),
    ("🪪", "Aadhaar KYC", "Digio API integration\nLand registration with OTP"),
    ("🔗", "Blockchain Audit", "Tamper-evident admin trail\nCustom BlockchainService.js"),
    ("🌐", "Multilingual UI", "Gujarati, Hindi & English\ni18next integration"),
]

cols = [(0.3, 1.35), (3.65, 1.35), (7.0, 1.35), (10.35, 1.35),
        (0.3, 4.0),  (3.65, 4.0),  (7.0, 4.0),  (10.35, 4.0)]

for idx, (icon, title, desc) in enumerate(features):
    l, t = cols[idx]
    b = box(sl, l, t, 2.9, 2.45, bg=ACCENT, border=GOLD)
    txt(sl, icon, l+0.1, t+0.1, 0.6, 0.55, 22, False, WHITE, PP_ALIGN.CENTER)
    txt(sl, title, l+0.7, t+0.12, 2.1, 0.45, 13, True, GOLD)
    txt(sl, desc,  l+0.1, t+0.65, 2.7, 1.6,  11, False, LGRAY)

# ════════════════════════════════════════════════════════
# SLIDE 4 — Tech Stack
# ════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_bg(sl)
header_bar(sl, "Technology Stack")
footer(sl, 4)

stacks = [
    ("🖥️ Frontend", [
        "React.js (Create React App)",
        "React Router v6",
        "Vanilla CSS",
        "i18next (Multilingual)",
        "Recharts (Analytics)",
        "Web Speech API (Voice Nav)",
    ]),
    ("⚙️ Backend", [
        "Node.js + Express.js",
        "MongoDB + Mongoose ODM",
        "JWT Authentication",
        "Passport.js (Google OAuth)",
        "Nodemailer (Gmail OTP)",
        "Helmet + Rate Limiting",
    ]),
    ("🤖 AI Service", [
        "Python FastAPI (Port 8005)",
        "TensorFlow / Keras",
        "EfficientNetB3 CNN Model",
        "OpenCV — CLAHE Enhancement",
        "Pillow (Fallback Processing)",
        "38 Crop Disease Classes",
    ]),
]

positions = [0.4, 4.65, 8.9]
for i, (heading, items) in enumerate(stacks):
    l = positions[i]
    box(sl, l, 1.3, 4.0, 5.5, bg=ACCENT, border=GOLD)
    txt(sl, heading, l+0.15, 1.4, 3.7, 0.55, 15, True, GOLD)
    box(sl, l+0.1, 1.9, 3.8, 0.03, bg=GOLD)
    for j, item in enumerate(items):
        txt(sl, "• " + item, l+0.2, 2.05 + j*0.75, 3.6, 0.65, 12.5, False, WHITE)

txt(sl, "Ports:  React  :3000   |   Express  :5000   |   FastAPI  :8005",
    0, 6.9, 13.33, 0.4, 12, True, GOLD, PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════
# SLIDE 5 — System Architecture
# ════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_bg(sl)
header_bar(sl, "System Architecture")
footer(sl, 5)

# React box
box(sl, 0.5, 1.35, 5.5, 1.1, bg=ACCENT, border=GOLD)
txt(sl, "🖥️  React Frontend  (Port 3000)", 0.7, 1.42, 5.1, 0.45, 13, True, GOLD)
txt(sl, "Login | Dashboards | AI UI | Cart | Market | KYC", 0.7, 1.8, 5.1, 0.4, 11, False, LGRAY)

# Arrow down
txt(sl, "▼  HTTP / JWT Bearer Token", 1.5, 2.55, 4, 0.4, 11, False, WHITE)

# Express box
box(sl, 0.5, 3.0, 5.5, 2.2, bg=ACCENT, border=GOLD)
txt(sl, "⚙️  Express.js Backend  (Port 5000)", 0.7, 3.07, 5.1, 0.45, 13, True, GOLD)
for j, r in enumerate(["Auth Routes — Signup / Login / OTP / Google OAuth",
                        "API Routes  — Land / Orders / Products / Market",
                        "Admin Routes — User Mgmt / Audit Logs",
                        "AI Proxy    — Forward images → FastAPI"]):
    txt(sl, "• " + r, 0.8, 3.58 + j*0.4, 5.0, 0.38, 11, False, LGRAY)

# Arrow right to FastAPI
txt(sl, "► image data", 6.15, 3.55, 1.5, 0.35, 10, False, WHITE)

# FastAPI box
box(sl, 7.7, 3.0, 5.0, 2.2, bg=ACCENT, border=GOLD)
txt(sl, "🤖  FastAPI AI  (Port 8005)", 7.9, 3.07, 4.6, 0.45, 13, True, GOLD)
for j, r in enumerate(["EfficientNetB3 CNN Model (.h5)",
                        "CLAHE Image Preprocessing",
                        "38 Disease Class Prediction",
                        "Returns: disease + treatment + confidence"]):
    txt(sl, "• " + r, 8.0, 3.58 + j*0.4, 4.5, 0.38, 11, False, LGRAY)

# Arrow down to MongoDB
txt(sl, "▼  Mongoose ODM", 1.5, 5.3, 4, 0.4, 11, False, WHITE)

# MongoDB box
box(sl, 0.5, 5.75, 12.2, 1.1, bg=RGBColor(0x0A,0x32,0x18), border=GOLD)
txt(sl, "🍃  MongoDB Atlas — Collections:", 0.7, 5.82, 4, 0.45, 13, True, GOLD)
txt(sl, "users  |  otps  |  products  |  orders  |  complaints  |  lands  |  diseasecases  |  auditlogs  |  marketrates",
    0.7, 6.25, 11.5, 0.45, 11, False, LGRAY)

# ════════════════════════════════════════════════════════
# SLIDE 6 — Authentication Flow
# ════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_bg(sl)
header_bar(sl, "Authentication & Security")
footer(sl, 6)

flow_steps = [
    "1.  User enters Username + Password  →  Login.js",
    "2.  POST /api/auth/login  →  server.js validates via bcrypt",
    "3.  generateOTP()  →  Otp.create()  →  Email sent via Nodemailer",
    "4.  User submits OTP  →  POST /api/auth/verify-otp",
    "5.  jwt.sign({ id, role }, secret, 7d)  →  JWT Token returned",
    "6.  Token stored in localStorage  →  User navigated to Dashboard",
]
for j, step in enumerate(flow_steps):
    box(sl, 0.5, 1.35 + j*0.82, 8.5, 0.7, bg=ACCENT, border=GOLD)
    txt(sl, step, 0.7, 1.42 + j*0.82, 8.2, 0.58, 12.5, False, WHITE)

# Security panel
box(sl, 9.3, 1.35, 3.7, 5.5, bg=ACCENT, border=GOLD)
txt(sl, "🔒 Security Layers", 9.5, 1.42, 3.3, 0.45, 13, True, GOLD)
for j, s in enumerate(["bcrypt (salt=10) password hash",
                        "JWT Token — 7 day expiry",
                        "OTP — 5 min expiry, deleted after use",
                        "Helmet.js — HTTP security headers",
                        "Rate Limit — 100 req / 15 min",
                        "express-validator input checks",
                        "Google OAuth (Passport.js)",
                        "Digio Aadhaar eKYC"]):
    txt(sl, "• " + s, 9.5, 1.95 + j*0.6, 3.4, 0.55, 11, False, LGRAY)

# ════════════════════════════════════════════════════════
# SLIDE 7 — AI Module
# ════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_bg(sl)
header_bar(sl, "AI Crop Disease Detection Module")
footer(sl, 7)

steps = [
    ("📷", "Farmer uploads crop image via CropDiagnostics.js"),
    ("🔧", "CLAHE preprocessing — adaptive contrast enhancement (OpenCV)"),
    ("📐", "Resize image to 300×300 — model input preparation"),
    ("🧠", "EfficientNetB3 CNN  →  model.predict()  →  38-class probabilities"),
    ("🎯", "Crop-specific filtering  +  Temperature Scaling (T = 0.2)"),
    ("✅", "Output: Disease name, confidence %, treatment, prevention, top-3"),
]
for j, (icon, label) in enumerate(steps):
    box(sl, 0.4, 1.35 + j*0.83, 7.5, 0.72, bg=ACCENT, border=GOLD)
    txt(sl, icon, 0.5,  1.42 + j*0.83, 0.7, 0.55, 18, False, WHITE, PP_ALIGN.CENTER)
    txt(sl, label, 1.25, 1.45 + j*0.83, 6.5, 0.55, 12.5, False, WHITE)

# Sample output card
box(sl, 8.2, 1.35, 4.7, 5.45, bg=ACCENT, border=GOLD)
txt(sl, "📋 Sample Output", 8.4, 1.42, 4.3, 0.45, 13, True, GOLD)
sample = [("Disease:", "Wheat Rust"),
          ("Confidence:", "94.2%"),
          ("Severity:", "High"),
          ("Treatment:", "Apply Propiconazole 25EC"),
          ("Prevention:", "Use resistant variety"),
          ("Image Quality:", "Blur: OK | Brightness: OK"),
          ("Top 3 Alt:", "Brown Spot, Blight, Healthy")]
for j, (k, v) in enumerate(sample):
    txt(sl, k, 8.4, 1.95 + j*0.65, 1.6, 0.55, 11, True, GOLD)
    txt(sl, v, 10.0, 1.95 + j*0.65, 2.8, 0.55, 11, False, WHITE)

txt(sl, "Model: crop_disease_model.h5  |  38 classes  |  Pillow fallback if OpenCV unavailable",
    0.4, 7.05, 12.5, 0.35, 10, False, LGRAY, PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════
# SLIDE 8 — Conclusion & Future Scope
# ════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_bg(sl)
header_bar(sl, "Conclusion & Future Scope")
footer(sl, 8)

box(sl, 0.4, 1.35, 6.0, 5.3, bg=ACCENT, border=GOLD)
txt(sl, "✅ What We Achieved", 0.6, 1.42, 5.6, 0.45, 14, True, GOLD)
for j, item in enumerate([
    "Full-stack MERN + Python AI web platform",
    "OTP + Google OAuth + JWT secure auth",
    "AI disease detection (90%+ accuracy target)",
    "Aadhaar KYC with Digio API integration",
    "E-Commerce with Cart, Checkout, Invoice",
    "Blockchain-based admin audit trail",
    "Complaint auto-routing + email notification",
    "Multilingual UI — Gujarati, Hindi, English",
]):
    txt(sl, "• " + item, 0.7, 1.97 + j*0.58, 5.5, 0.52, 12, False, WHITE)

box(sl, 6.9, 1.35, 6.0, 5.3, bg=ACCENT, border=GOLD)
txt(sl, "🚀 Future Scope", 7.1, 1.42, 5.6, 0.45, 14, True, GOLD)
for j, item in enumerate([
    "React Native mobile app for farmers",
    "Real-time IoT sensor data integration",
    "Live government scheme API connection",
    "SMS OTP via Twilio",
    "Drone imagery for large field scanning",
    "Crop price prediction using ML regression",
    "Online payment gateway (Razorpay/Stripe)",
    "Progressive Web App (PWA) support",
]):
    txt(sl, "• " + item, 7.1, 1.97 + j*0.58, 5.6, 0.52, 12, False, WHITE)

# ════════════════════════════════════════════════════════
# SLIDE 9 — Thank You
# ════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_bg(sl)
box(sl, 0, 0, 13.33, 0.5, bg=RGBColor(0x0A,0x20,0x10))

crest2 = sl.shapes.add_shape(9, Inches(5.7), Inches(0.65), Inches(1.9), Inches(1.9))
crest2.fill.solid(); crest2.fill.fore_color.rgb = GOLD
crest2.line.color.rgb = WHITE; crest2.line.width = Pt(2)
txt(sl, "DDU", 5.7, 1.05, 1.9, 0.7, 22, True, BG, PP_ALIGN.CENTER)

txt(sl, "THANK YOU", 0, 2.7, 13.33, 0.85, 44, True, WHITE, PP_ALIGN.CENTER)
txt(sl, "Khedut Bandhu – Smart Agriculture Management System (SAMS)",
    0, 3.55, 13.33, 0.5, 15, True, GOLD, PP_ALIGN.CENTER)

box(sl, 2.0, 4.15, 9.33, 0.04, bg=GOLD)

txt(sl, "Presented By:", 1.5, 4.3, 4.5, 0.4, 13, True, GOLD)
txt(sl, "Padaliya Princekumar R.  (IT-085)\nPatel Ansh K.  (IT-091)\nB.Tech IT, Semester VI",
    1.5, 4.7, 5.5, 1.0, 13, False, WHITE)

txt(sl, "Guided By:", 7.5, 4.3, 4.5, 0.4, 13, True, GOLD)
txt(sl, "Prof. Sunil K. Vithlani\nDept. of Information Technology\nDharmsinh Desai University",
    7.5, 4.7, 5.2, 1.0, 13, False, WHITE)

box(sl, 2.0, 5.8, 9.33, 0.04, bg=GOLD)
txt(sl, "Dharmsinh Desai University, Nadiad – 387001, Gujarat  |  April 2026",
    0, 5.9, 13.33, 0.4, 12, False, LGRAY, PP_ALIGN.CENTER)
txt(sl, "🔗 github.com/apk3206/Khedut-Bandhu",
    0, 6.35, 13.33, 0.4, 12, False, GOLD, PP_ALIGN.CENTER)

# ── Save ─────────────────────────────────────────────────
out = r"d:\frontend\Khedut_Bandhu_Presentation.pptx"
prs.save(out)
print(f"✅ PPT saved: {out}")
print(f"   Slides: {len(prs.slides)}")
